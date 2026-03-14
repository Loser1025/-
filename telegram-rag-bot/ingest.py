"""
PDF/PPTXファイルをChromaDBに取り込むスクリプト
使い方: python ingest.py
"""

import os
import glob
import fitz  # pymupdf
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
import chromadb
from chromadb.utils import embedding_functions
from dotenv import load_dotenv

load_dotenv()

PDF_DIR = "./pdfs"
CHROMA_DIR = "./chroma_db"
CHUNK_SIZE = 500      # 1チャンクの文字数
CHUNK_OVERLAP = 50    # チャンク間の重複文字数

def extract_text_from_pdf(pdf_path: str) -> str:
    """PDFからテキストを抽出する（テキストなければOCR）"""
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()

    # テキストが取れなかった場合はOCRで読み取る
    if not text.strip():
        print("  → テキストなし、OCRで読み取り中...")
        images = convert_from_path(pdf_path, dpi=200)
        for i, img in enumerate(images, 1):
            page_text = pytesseract.image_to_string(img, lang="jpn")
            if page_text.strip():
                text += f"【ページ{i}】\n{page_text}\n"
            if i % 5 == 0:
                print(f"  → {i}/{len(images)}ページ完了")

    return text

def extract_text_from_pptx(pptx_path: str) -> str:
    """PPTXからテキストを抽出する"""
    prs = Presentation(pptx_path)
    text = ""
    for i, slide in enumerate(prs.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_text += line + "\n"
        if slide_text:
            text += f"【スライド{i}】\n{slide_text}\n"
    return text

def split_into_chunks(text: str, filename: str) -> list[dict]:
    """テキストをチャンクに分割する"""
    chunks = []
    start = 0
    chunk_index = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk_text = text[start:end].strip()
        if chunk_text:
            chunks.append({
                "text": chunk_text,
                "id": f"{filename}_{chunk_index}",
                "metadata": {"source": filename, "chunk": chunk_index}
            })
            chunk_index += 1
        start = end - CHUNK_OVERLAP
    return chunks

def main():
    # 日本語対応の埋め込みモデル
    print("埋め込みモデルを読み込み中... (初回は数分かかります)")
    ef = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name="intfloat/multilingual-e5-small"
    )

    # ChromaDB初期化
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    collection = client.get_or_create_collection(
        name="knowledge",
        embedding_function=ef
    )

    # PDF・PPTXファイルを処理
    pdf_files = glob.glob(os.path.join(PDF_DIR, "*.pdf"))
    pptx_files = glob.glob(os.path.join(PDF_DIR, "*.pptx"))
    all_files = pdf_files + pptx_files

    if not all_files:
        print(f"⚠️  {PDF_DIR} にPDF/PPTXファイルが見つかりません")
        return

    print(f"{len(all_files)}件のファイルを処理します（PDF: {len(pdf_files)}件、PPTX: {len(pptx_files)}件）")

    for file_path in all_files:
        filename = os.path.basename(file_path)
        print(f"処理中: {filename}")

        # テキスト抽出
        if filename.endswith(".pptx"):
            text = extract_text_from_pptx(file_path)
        else:
            text = extract_text_from_pdf(file_path)
        if not text.strip():
            print(f"  ⚠️  テキストが抽出できませんでした（画像PDFの可能性）")
            continue

        print(f"  抽出文字数: {len(text):,}文字")

        # チャンク分割
        chunks = split_into_chunks(text, filename)
        print(f"  チャンク数: {len(chunks)}")

        # ChromaDBに保存（既存は上書き）
        existing = collection.get(where={"source": filename})
        if existing["ids"]:
            collection.delete(where={"source": filename})

        collection.add(
            documents=[c["text"] for c in chunks],
            ids=[c["id"] for c in chunks],
            metadatas=[c["metadata"] for c in chunks]
        )
        print(f"  ✅ 保存完了")

    total = collection.count()
    print(f"\n✅ 完了！ 合計チャンク数: {total}")

if __name__ == "__main__":
    main()
